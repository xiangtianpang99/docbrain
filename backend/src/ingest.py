import os
import glob
from typing import List, Optional
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

# Loaders
import pypdf
import docx
import openpyxl
from pptx import Presentation
from src.config_manager import config_manager

class IngestionEngine:
    def __init__(self, persist_directory: str = None, model_name: str = "all-MiniLM-L6-v2"):
        if persist_directory is None:
            # Resolve to absolute path relative to project root
            # Resolve to absolute path relative to project root (docbrain root)
            # ingest.py is in backend/src, so:
            # os.path.dirname(os.path.abspath(__file__)) -> backend/src
            # os.path.dirname(...) -> backend
            # os.path.dirname(...) -> docbrain (Project Root)
            backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            project_root = os.path.dirname(backend_dir)
            persist_directory = os.path.join(project_root, "chroma_db")
            
        self.persist_directory = persist_directory

        # Check for local model
        # Models are in backend/models
        backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        local_model_path = os.path.join(backend_dir, "models", model_name)
        if os.path.exists(local_model_path):
            print(f"Found local model at: {local_model_path}")
            model_name = local_model_path
        
        print(f"Loading embedding model: {model_name}...")
        # Force CPU to avoid iGPU spikes on some Windows configs
        model_kwargs = {'device': 'cpu'}
        encode_kwargs = {'normalize_embeddings': False}
        
        self.embedding_model = HuggingFaceEmbeddings(
            model_name=model_name,
            model_kwargs=model_kwargs,
            encode_kwargs=encode_kwargs
        )
        self.vector_store = Chroma(
            persist_directory=self.persist_directory,
            embedding_function=self.embedding_model
        )
        self.busy_jobs = 0
        import time
        self.last_update_time = time.time()

    def start_job(self):
        self.busy_jobs += 1

    def end_job(self):
        if self.busy_jobs > 0:
             self.busy_jobs -= 1
        import time
        self.last_update_time = time.time()

    def parse_file(self, file_path: str) -> List[Document]:
        """
        根据文件扩展名解析单个文件并返回 Document 对象列表。
        """
        ext = os.path.splitext(file_path)[1].lower()
        content = ""
        try:
            if ext in [".txt", ".md"]:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            elif ext == ".pdf":
                reader = pypdf.PdfReader(file_path)
                text_list = []
                for page in reader.pages:
                    text_list.append(page.extract_text())
                content = "\n".join(text_list)
            elif ext == ".docx":
                doc = docx.Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".doc":
                # Using unstructured for legacy .doc
                from unstructured.partition.auto import partition
                elements = partition(filename=file_path)
                content = "\n".join([str(el) for el in elements])
            elif ext == ".xlsx":
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text_list = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        row_text = " ".join([str(cell) for cell in row if cell is not None])
                        if row_text:
                            text_list.append(row_text)
                content = "\n".join(text_list)
            elif ext == ".xls":
                # Using pandas with xlrd for legacy .xls
                import pandas as pd
                df_dict = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
                text_list = []
                for sheet_name, df in df_dict.items():
                    text_list.append(f"Sheet: {sheet_name}")
                    text_list.append(df.to_string(index=False))
                content = "\n".join(text_list)
            elif ext == ".pptx":
                prs = Presentation(file_path)
                text_list = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_list.append(shape.text)
                content = "\n".join(text_list)
            elif ext == ".ppt":
                # Using unstructured for legacy .ppt
                from unstructured.partition.auto import partition
                elements = partition(filename=file_path)
                content = "\n".join([str(el) for el in elements])
            else:
                print(f"Unsupported file format: {file_path}")
                return []

            if not content.strip():
                return []
            
            # Get file stats for metadata
            stats = os.stat(file_path)
            metadata = {
                "source": file_path,
                "title": os.path.basename(file_path),
                "type": "file",
                "extension": os.path.splitext(file_path)[1].lower(),
                "file_size": stats.st_size,
                "mtime": stats.st_mtime
            }
            
            # Return as a single document (splitter will handle chunking)
            return [Document(page_content=content, metadata=metadata)]

        except Exception as e:
            print(f"Error parsing {file_path}: {e}")
            return []

    def split_documents(self, documents: List[Document]) -> List[Document]:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=300,
            separators=["\n\n", "\n", " ", ""]
        )
        return text_splitter.split_documents(documents)

    def process_file(self, file_path: str, additional_duration: int = 0):
        """
        索引单个文件: 删除旧向量 -> 解析 -> 切分 -> 添加新向量。
        additional_duration: 需要添加到现有文件时长的秒数。
        """
        self.start_job()
        try:
            abs_path = os.path.abspath(file_path)
            print(f"正在处理文件: {abs_path}")
            
            # 1.获取现有时长 (如果有)
            existing_duration = 0
            try:
                results = self.vector_store._collection.get(where={"source": abs_path})
                if results and results.get("metadatas"):
                    # All chunks for one file should have the same duration
                    existing_duration = results["metadatas"][0].get("duration", 0)
            except Exception:
                pass

            total_duration = int(existing_duration + additional_duration)

            # 2. 删除此文件的现有向量以避免重复
            try:
                self.vector_store._collection.delete(where={"source": abs_path})
                # 清理潜在的旧路径格式
                rel_path = os.path.relpath(abs_path)
                if rel_path != abs_path:
                     self.vector_store._collection.delete(where={"source": rel_path})
            except Exception:
                pass

            # 3. Parse
            documents = self.parse_file(abs_path)
            if not documents:
                return

            # Inject total duration into metadata
            for doc in documents:
                doc.metadata["duration"] = total_duration

            # 4. 切分 (Split)
            chunks = self.split_documents(documents)
            
            # 5. 添加 (Add)
            if chunks:
                self.vector_store.add_documents(chunks)
                self.vector_store.persist()
                print(f"已索引 {len(chunks)} 个分块。总时长: {total_duration}秒")
        finally:
            self.end_job()
        
    def remove_document(self, file_path: str):
        """
        Remove vectors associated with a file.
        """
        self.start_job()
        try:
            file_path = os.path.abspath(file_path)
            print(f"Removing documents for: {file_path}")
            self.vector_store._collection.delete(where={"source": file_path})
            self.vector_store.persist()
        except Exception as e:
            print(f"Error removing {file_path}: {e}")
        finally:
            self.end_job()

    def remove_documents_by_root(self, root_path: str):
        """
        移除属于特定根目录的所有文档。
        由于 Chroma 不支持方便的 'startswith' 删除查询，
        所以我们获取所有元数据，在 Python 中过滤，然后按 ID 删除。
        """
        self.start_job()
        try:
            root_path = os.path.abspath(root_path)
            print(f"Cleaning up documents from root: {root_path}")
            
            # 1. Get all documents
            result = self.vector_store._collection.get()
            if not result or not result['ids']:
                return

            ids_to_delete = []
            sources = result['metadatas']
            ids = result['ids']

            # 2. Filter match
            for i, metadata in enumerate(sources):
                if metadata and 'source' in metadata:
                    source = metadata['source']
                    # Check if source starts with root_path
                    # using commonpath is safer for path logic
                    try:
                        if os.path.commonpath([root_path, source]) == root_path:
                            ids_to_delete.append(ids[i])
                    except ValueError:
                        continue # Different drivers

            # 3. Delete
            if ids_to_delete:
                print(f"Found {len(ids_to_delete)} chunks to remove.")
                self.vector_store._collection.delete(ids=ids_to_delete)
                self.vector_store.persist()
                print("Cleanup complete.")
            else:
                print("No documents found for this root.")

        except Exception as e:
            print(f"Error cleaning root {root_path}: {e}")
        finally:
            self.end_job()

    def ingest_webpage(self, url: str, title: str, content: str, additional_duration: int = 0):
        """
        索引网页内容。
        additional_duration: 在此页面上花费的秒数，累加到总数。
        """
        self.start_job()
        try:
            print(f"正在索引网页: {title} ({url})")
            
            # 1. Fetch existing duration if any
            existing_duration = 0
            try:
                results = self.vector_store._collection.get(where={"source": url})
                if results and results.get("metadatas"):
                    existing_duration = results["metadatas"][0].get("duration", 0)
            except Exception:
                pass

            total_duration = int(existing_duration + additional_duration)

            # 2. Deduplication: Remove old entry for this URL
            try:
                self.vector_store._collection.delete(where={"source": url})
                self.vector_store.persist()
            except Exception:
                pass
                
            # 3. Prepare metadata
            import time
            metadata = {
                "source": url,
                "title": title,
                "type": "webpage",
                "extension": ".html",
                "duration": total_duration,
                "mtime": time.time()
            }
            
            # 4. Create document and split
            doc = Document(page_content=content, metadata=metadata)
            chunks = self.split_documents([doc])
            
            # 5. Add to store
            if chunks:
                self.vector_store.add_documents(chunks)
                self.vector_store.persist()
                print(f"Webpage indexed: {len(chunks)} chunks. Total duration: {total_duration}s")
                return len(chunks)
            return 0
        finally:
            self.end_job()

    def ingest_directory(self, source_dir: str):
        """
        索引目录中的所有支持文件，跳过系统和临时文件夹。
        """
        ignore_dirs = {
            "node_modules", ".git", ".venv", ".vscode", "__pycache__", 
            "System Volume Information", "$RECYCLE.BIN", ".idea", ".DS_Store",
            "venv", "env", "tmp", "temp"
        }
        
        if not os.path.isabs(source_dir):
            # If relative, valid against process CWD or project root
            # Priority: 1. CWD (if exists) 2. Backend Root (if exists)
            if not os.path.exists(source_dir):
                backend_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                potential_path = os.path.join(backend_root, source_dir)
                if os.path.exists(potential_path):
                    source_dir = potential_path

        print(f"Scanning directory: {source_dir}")
        patterns = [
            "**/*.txt", "**/*.md", "**/*.pdf", 
            "**/*.docx", "**/*.doc", 
            "**/*.xlsx", "**/*.xls", 
            "**/*.pptx", "**/*.ppt"
        ]
        all_files = []
        
        # Using os.walk for better control over directory exclusion
        for root, dirs, files in os.walk(source_dir):
            # In-place modification of dirs to skip ignored directories
            dirs[:] = [d for d in dirs if d not in ignore_dirs and not d.startswith(".")]
            
            for file in files:
                if any(file.lower().endswith(p.replace("**/*", "")) for p in patterns):
                    all_files.append(os.path.join(root, file))
        
        print(f"Found {len(all_files)} files to process.")
        for file_path in all_files:
            # process_file will handle abspath conversion
            self.process_file(file_path)
        
        print("Ingestion complete.")

if __name__ == "__main__":
    ingestor = IngestionEngine()
    ingestor.ingest_directory("./data")
